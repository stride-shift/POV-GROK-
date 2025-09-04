import os
import docx
import PyPDF2
import pandas as pd
from pptx import Presentation
import markdown
import chardet

class FileReader:
    """A comprehensive file reader that supports multiple file formats."""
    
    @staticmethod
    def detect_encoding(file_path):
        """Detect the encoding of a file."""
        with open(file_path, 'rb') as file:
            raw_data = file.read()
            result = chardet.detect(raw_data)
            return result['encoding']

    @staticmethod
    def read_docx(file_path):
        """
        Read content from a .docx file.
        Returns a dictionary containing:
        - full_text: The complete text content
        - paragraphs: List of paragraphs
        - tables: List of tables
        """
        try:
            doc = docx.Document(file_path)
            
            # Extract text from paragraphs
            paragraphs = [para.text for para in doc.paragraphs]
            
            # Extract tables
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)
            
            return {
                'full_text': '\n'.join(paragraphs),
                'paragraphs': paragraphs,
                'tables': tables
            }
        except Exception as e:
            raise Exception(f"Error reading DOCX file: {str(e)}")

    @staticmethod
    def read_pdf(file_path):
        """
        Read content from a PDF file.
        Returns a dictionary containing:
        - full_text: The complete text content
        - pages: List of text content by page
        - metadata: PDF metadata
        """
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                # Extract text from each page
                pages = []
                for page in pdf_reader.pages:
                    pages.append(page.extract_text())
                
                # Get metadata
                metadata = pdf_reader.metadata if pdf_reader.metadata else {}
                
                return {
                    'full_text': '\n'.join(pages),
                    'pages': pages,
                    'metadata': metadata
                }
        except Exception as e:
            raise Exception(f"Error reading PDF file: {str(e)}")

    @staticmethod
    def read_pptx(file_path):
        """
        Read content from a PowerPoint file.
        Returns a dictionary containing:
        - slides: List of slides with their content
        - notes: List of slide notes
        - shapes: List of shapes/textboxes per slide
        """
        try:
            prs = Presentation(file_path)
            
            slides = []
            notes = []
            shapes = []
            
            for slide in prs.slides:
                # Get slide text
                slide_text = []
                slide_shapes = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                        slide_shapes.append({
                            'type': shape.shape_type,
                            'text': shape.text
                        })
                
                slides.append('\n'.join(slide_text))
                shapes.append(slide_shapes)
                
                # Get notes
                if slide.has_notes_slide:
                    note_text = slide.notes_slide.notes_text_frame.text
                    notes.append(note_text)
                else:
                    notes.append('')
            
            return {
                'slides': slides,
                'notes': notes,
                'shapes': shapes
            }
        except Exception as e:
            raise Exception(f"Error reading PPTX file: {str(e)}")

    @staticmethod
    def read_xlsx(file_path):
        """
        Read content from an Excel file.
        Returns a dictionary containing:
        - sheets: Dictionary of dataframes for each sheet
        - sheet_names: List of sheet names
        """
        try:
            # Read all sheets
            xlsx = pd.ExcelFile(file_path)
            sheets = {}
            
            for sheet_name in xlsx.sheet_names:
                df = pd.read_excel(xlsx, sheet_name)
                sheets[sheet_name] = df
            
            return {
                'sheets': sheets,
                'sheet_names': xlsx.sheet_names
            }
        except Exception as e:
            raise Exception(f"Error reading XLSX file: {str(e)}")

    @staticmethod
    def read_txt(file_path):
        """
        Read content from a text file.
        Returns a dictionary containing:
        - content: The complete text content
        - lines: List of lines
        - encoding: Detected encoding
        """
        try:
            encoding = FileReader.detect_encoding(file_path)
            
            with open(file_path, 'r', encoding=encoding) as file:
                content = file.read()
                lines = content.split('\n')
                
            return {
                'content': content,
                'lines': lines,
                'encoding': encoding
            }
        except Exception as e:
            raise Exception(f"Error reading TXT file: {str(e)}")

    @staticmethod
    def read_md(file_path):
        """
        Read content from a Markdown file.
        Returns a dictionary containing:
        - raw_content: The raw markdown content
        - html_content: Converted HTML content
        - lines: List of lines
        """
        try:
            encoding = FileReader.detect_encoding(file_path)
            
            with open(file_path, 'r', encoding=encoding) as file:
                content = file.read()
                lines = content.split('\n')
                
                # Convert to HTML
                html = markdown.markdown(content)
                
            return {
                'raw_content': content,
                'html_content': html,
                'lines': lines
            }
        except Exception as e:
            raise Exception(f"Error reading MD file: {str(e)}")

    @classmethod
    def read_file(cls, file_path):
        """
        Main method to read any supported file type.
        Automatically detects file type and uses appropriate reader.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_extension = os.path.splitext(file_path)[1].lower()
        
        # Map file extensions to reader methods
        readers = {
            '.docx': cls.read_docx,
            '.pdf': cls.read_pdf,
            '.pptx': cls.read_pptx,
            '.xlsx': cls.read_xlsx,
            '.txt': cls.read_txt,
            '.md': cls.read_md
        }
        
        if file_extension not in readers:
            raise ValueError(f"Unsupported file type: {file_extension}")
        
        return readers[file_extension](file_path)

# Example usage
if __name__ == "__main__":
    # Example usage for each file type
    try:
        # Read DOCX
        docx_result = FileReader.read_file("example.docx")
        print("DOCX Content:", docx_result['full_text'][:100])
        
        # Read PDF
        pdf_result = FileReader.read_file("example.pdf")
        print("PDF Content:", pdf_result['full_text'][:100])
        
        # Read PPTX
        pptx_result = FileReader.read_file("example.pptx")
        print("PPTX Slides:", len(pptx_result['slides']))
        
        # Read XLSX
        xlsx_result = FileReader.read_file("example.xlsx")
        print("Excel Sheets:", xlsx_result['sheet_names'])
        
        # Read TXT
        txt_result = FileReader.read_file("example.txt")
        print("TXT Content:", txt_result['content'][:100])
        
        # Read MD
        md_result = FileReader.read_file("example.md")
        print("MD Content:", md_result['raw_content'][:100])
        
    except Exception as e:
        print(f"Error: {str(e)}") 